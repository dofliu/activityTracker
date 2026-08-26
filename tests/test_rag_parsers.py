import os
import pytest
from rag.parsers.parser_hub import parser_hub
from rag.parsers.text_parser import TextParser
from rag.parsers.office_parser import DocxParser, PptxParser, ExcelParser
from rag.parsers.pdf_parser import PdfParser


def test_text_parser(tmp_path):
    f = tmp_path / "sample.md"
    f.write_text("# Title\nThis is a sample markdown document for RAG testing.", encoding="utf-8")
    doc = parser_hub.parse_file(str(f))
    assert doc.filename == "sample.md"
    assert doc.file_type == "md"
    assert len(doc.sections) >= 1
    assert "This is a sample markdown document" in doc.sections[0].content


def test_docx_parser(tmp_path):
    import docx
    f = tmp_path / "test.docx"
    d = docx.Document()
    d.add_heading("Section 1 Header", level=1)
    d.add_paragraph("Paragraph inside section 1.")
    d.save(str(f))

    doc = parser_hub.parse_file(str(f))
    assert doc.filename == "test.docx"
    assert doc.file_type == "docx"
    assert len(doc.sections) >= 1
    assert any("Paragraph inside section 1" in s.content for s in doc.sections)


def test_excel_parser(tmp_path):
    import openpyxl
    f = tmp_path / "data.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "SheetTest"
    ws.append(["ID", "Name", "Score"])
    ws.append([1, "Alice", 95])
    wb.save(str(f))

    doc = parser_hub.parse_file(str(f))
    assert doc.filename == "data.xlsx"
    assert doc.file_type == "xlsx"
    assert len(doc.sections) >= 1
    assert "Alice" in doc.sections[0].content
    assert doc.sections[0].sheet_name == "SheetTest"


def test_pptx_parser(tmp_path):
    from pptx import Presentation
    f = tmp_path / "presentation.pptx"
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "DeskRAG Overview"
    prs.save(str(f))

    doc = parser_hub.parse_file(str(f))
    assert doc.filename == "presentation.pptx"
    assert doc.file_type == "pptx"
    assert len(doc.sections) >= 1
    assert "DeskRAG Overview" in doc.sections[0].content
    assert doc.sections[0].slide_number == 1
