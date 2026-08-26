from pathlib import Path
from typing import List
from rag.parsers.base import BaseParser, ParsedDocument, ParsedSection


class DocxParser(BaseParser):
    def parse(self, file_path: str) -> ParsedDocument:
        path_obj = Path(file_path)
        sections: List[ParsedSection] = []

        try:
            import docx
            doc = docx.Document(file_path)
            current_heading = ""
            current_paragraphs = []

            for p in doc.paragraphs:
                text = p.text.strip()
                if not text:
                    continue
                if p.style and p.style.name and p.style.name.startswith("Heading"):
                    if current_paragraphs:
                        sections.append(ParsedSection(
                            content="\n".join(current_paragraphs),
                            section_title=current_heading or "Introduction"
                        ))
                        current_paragraphs = []
                    current_heading = text
                else:
                    current_paragraphs.append(text)

            if current_paragraphs:
                sections.append(ParsedSection(
                    content="\n".join(current_paragraphs),
                    section_title=current_heading or "Body"
                ))

            # Also extract tables
            for idx, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    table_data.append(" | ".join(row_cells))
                if table_data:
                    sections.append(ParsedSection(
                        content="\n".join(table_data),
                        section_title=f"Table {idx+1}"
                    ))
        except Exception as e:
            sections.append(ParsedSection(content=f"[Docx Parse Error: {str(e)}]"))

        return ParsedDocument(
            file_path=file_path,
            filename=path_obj.name,
            file_type="docx",
            sections=sections
        )


class PptxParser(BaseParser):
    def parse(self, file_path: str) -> ParsedDocument:
        path_obj = Path(file_path)
        sections: List[ParsedSection] = []

        try:
            from pptx import Presentation
            prs = Presentation(file_path)

            for idx, slide in enumerate(prs.slides):
                slide_num = idx + 1
                slide_texts = []
                slide_title = f"Slide {slide_num}"

                # Check for slide title
                if slide.shapes.title and slide.shapes.title.text:
                    slide_title = f"Slide {slide_num}: {slide.shapes.title.text.strip()}"

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            txt = paragraph.text.strip()
                            if txt and txt not in slide_texts:
                                slide_texts.append(txt)

                # Check speaker notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_texts.append(f"[Speaker Notes: {notes}]")

                if slide_texts:
                    sections.append(ParsedSection(
                        content="\n".join(slide_texts),
                        slide_number=slide_num,
                        section_title=slide_title
                    ))
        except Exception as e:
            sections.append(ParsedSection(content=f"[Pptx Parse Error: {str(e)}]"))

        return ParsedDocument(
            file_path=file_path,
            filename=path_obj.name,
            file_type="pptx",
            sections=sections
        )


class ExcelParser(BaseParser):
    def parse(self, file_path: str) -> ParsedDocument:
        path_obj = Path(file_path)
        sections: List[ParsedSection] = []

        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                rows_text = []

                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None and str(cell).strip() != "" for cell in row):
                        row_str = " | ".join(str(cell).strip() if cell is not None else "" for cell in row)
                        rows_text.append(row_str)

                if rows_text:
                    sections.append(ParsedSection(
                        content="\n".join(rows_text[:500]),
                        sheet_name=sheet_name,
                        section_title=f"Sheet: {sheet_name}"
                    ))
            wb.close()
        except Exception as e:
            sections.append(ParsedSection(content=f"[Excel Parse Error: {str(e)}]"))

        return ParsedDocument(
            file_path=file_path,
            filename=path_obj.name,
            file_type="xlsx",
            sections=sections
        )
